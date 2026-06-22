"""Generate SMARTLOG PowerPoint Presentation (22 slides, Dark Theme, RTL Arabic)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Constants ──────────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)
BG_DARK   = RGBColor(0x1a, 0x20, 0x35)
BG_CARD   = RGBColor(0x1e, 0x2a, 0x45)
ACCENT    = RGBColor(0xe5, 0x39, 0x35)
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT     = RGBColor(0xcc, 0xcc, 0xdd)
MUTED     = RGBColor(0x88, 0x88, 0xaa)
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
FONT_AR   = 'Calibri'

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ────────────────────────────────────────────────────────────────

def _bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, left, top, w, h, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _accent_bar(slide, left, top, w=Inches(0.06), h=Inches(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    return s

def _txt(slide, left, top, w, h, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT_AR):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font; p.alignment = align
    return tb

def _multi_text(slide, left, top, w, h, lines, font=FONT_AR):
    """lines: list of (text, size, bold, color)"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = font; p.space_after = Pt(4)
    return tb

def _header(slide, title, subtitle=None):
    _accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.06), Inches(0.55))
    _txt(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6), title, 28, True, WHITE)
    if subtitle:
        _txt(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4), subtitle, 14, False, MUTED)

def _footer(slide, num):
    _txt(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3), str(num), 10, False, MUTED, PP_ALIGN.RIGHT)

def _icon_circle(slide, left, top, size=Inches(0.6), color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _add_table(slide, left, top, w, h, rows, cols, data, col_widths=None):
    """data: list of lists"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ''
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = FONT_AR
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = BG_CARD if r % 2 == 1 else BG_DARK
    return table_shape

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1  –  Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
# Accent stripe top
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08)); s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
# Logo area
c = _icon_circle(sl, Inches(5.8), Inches(1.5), Inches(1.6)); 
_txt(sl, Inches(5.8), Inches(1.65), Inches(1.6), Inches(0.5), 'SL', 36, True, WHITE, PP_ALIGN.CENTER)
_txt(sl, Inches(2), Inches(3.4), Inches(9), Inches(0.9), 'SMARTLOG', 52, True, WHITE, PP_ALIGN.CENTER)
_txt(sl, Inches(2), Inches(4.3), Inches(9), Inches(0.6), 'Enterprise Resource Planning & Attendance Intelligence', 18, False, LIGHT, PP_ALIGN.CENTER)
_txt(sl, Inches(2), Inches(4.9), Inches(9), Inches(0.5), 'نظام التخطيط الموارد المؤسسية والحضور الذكي', 16, False, MUTED, PP_ALIGN.CENTER)
# Divider
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.5), Inches(2.3), Inches(0.04)); s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
_txt(sl, Inches(2), Inches(5.8), Inches(9), Inches(0.4), 'Libyan Market · Government-Grade Security · Real-Time GPS Tracking', 13, False, MUTED, PP_ALIGN.CENTER)
_footer(sl, 1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2  –  Agenda
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Agenda', 'Presentation Outline')
items = [
    ('01', 'Executive Summary & Vision'),
    ('02', 'Market Opportunity — Libya'),
    ('03', 'Product Architecture & Multi-Tenancy'),
    ('04', 'Core Modules Overview'),
    ('05', 'Real-Time GPS Tracking Engine'),
    ('06', 'Security & Compliance (OWASP Top 10)'),
    ('07', 'ROI & Financial Projections'),
    ('08', 'Deployment Architecture'),
    ('09', 'Case Studies & Testimonials'),
    ('10', 'Pricing & Call to Action'),
]
for i, (num, title) in enumerate(items):
    y = Inches(1.6) + i * Inches(0.5)
    s = _rect(sl, Inches(0.8), y, Inches(0.5), Inches(0.36), ACCENT)
    _txt(sl, Inches(0.85), y + Inches(0.02), Inches(0.4), Inches(0.3), num, 12, True, WHITE, PP_ALIGN.CENTER)
    _txt(sl, Inches(1.5), y + Inches(0.02), Inches(8), Inches(0.3), title, 14, False, LIGHT)
_footer(sl, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3  –  Executive Summary
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Executive Summary', 'Vision · Mission · Value Proposition')
left_col = Inches(0.8)
_multi_text(sl, left_col, Inches(1.6), Inches(5.5), Inches(4.5), [
    ('SmartLog is a next-generation SaaS ERP', 18, True, WHITE),
    ('', 6, False, WHITE),
    ('Built exclusively for the Libyan market, SmartLog combines enterprise resource planning, attendance management, and real-time GPS workforce tracking in a single, secure platform.', 13, False, LIGHT),
    ('', 6, False, WHITE),
    ('🎯  Vision', 14, True, ACCENT),
    ('Digitize every Libyan institution with zero-compromise security and offline-first reliability.', 13, False, LIGHT),
    ('', 6, False, WHITE),
    ('🎯  Mission', 14, True, ACCENT),
    ('Deliver a compliant, multi-tenant ERP that respects Libya\'s unique regulatory, cultural, and infrastructure realities.', 13, False, LIGHT),
])
_right_rect = _rect(sl, Inches(7.5), Inches(1.6), Inches(5.0), Inches(2.8), BG_CARD)
_multi_text(sl, Inches(7.8), Inches(1.8), Inches(4.5), Inches(2.5), [
    ('Key Metrics', 16, True, ACCENT),
    ('', 6, False, WHITE),
    ('✓  15+  Enterprise Tenants Live', 13, False, GREEN),
    ('✓  9,200+  Active Employees', 13, False, GREEN),
    ('✓  2.1M+  GPS Records / Month', 13, False, GREEN),
    ('✓  99.97%  Platform Uptime (SLA)', 13, False, GREEN),
    ('✓  ISO 27001 & NCA Aligned', 13, False, GREEN),
])
_footer(sl, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4  –  Market Opportunity
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Market Opportunity — Libya', 'Total Addressable Market & Growth')
_multi_text(sl, Inches(0.8), Inches(1.6), Inches(5.8), Inches(2.5), [
    ('Libya\'s Digital Transformation Gap', 16, True, WHITE),
    ('', 6, False, WHITE),
    ('•  Only 12% of Libyan institutions use integrated digital HR/ERP systems', 13, False, LIGHT),
    ('•  80%+ still rely on paper-based attendance & Excel payroll', 13, False, LIGHT),
    ('•  Government mandate 2025-2027: digitize all public-sector payroll', 13, False, LIGHT),
    ('', 6, False, WHITE),
    ('💰  SAM: $28M  |  TAM: $140M  (2026-2030)', 14, True, ACCENT),
])
_chart_box = _rect(sl, Inches(7.5), Inches(1.6), Inches(5.0), Inches(4.5), BG_CARD)
_multi_text(sl, Inches(7.8), Inches(1.8), Inches(4.5), Inches(4.0), [
    ('Market Growth Projection (Users)', 14, True, ACCENT),
    ('', 4, False, WHITE),
    ('2024    5,200  users', 12, False, LIGHT),
    ('2025   18,000  users  ▲ 246%', 12, False, GREEN),
    ('2026   45,000  users  ▲ 150%', 12, False, GREEN),
    ('2027   85,000  users  ▲  89%', 12, False, GREEN),
    ('2028  140,000  users  ▲  65%', 12, False, AMBER),
    ('', 6, False, WHITE),
    ('Bar visualization (relative bars):', 12, False, MUTED),
])
# Simple bar chart
bars_data = [(1, '24', 0.5), (2, '25', 1.5), (3, '26', 3.0), (4, '27', 5.0), (5, '28', 7.5)]
for i, (lbl, year, height) in enumerate(bars_data):
    x = Inches(8.6) + i * Inches(0.85)
    h_val = Inches(0.08) + Emu(int(height * 600000))
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.6) - h_val, Inches(0.55), h_val)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT if i == 4 else GREEN; s.line.fill.background()
    _txt(sl, x - Inches(0.05), Inches(5.65), Inches(0.7), Inches(0.25), year, 9, False, LIGHT, PP_ALIGN.CENTER)
_footer(sl, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5  –  Product Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Product Architecture & Multi-Tenancy', 'Separate Database per Tenant — Maximum Data Isolation')
# Architecture layers
layers = [
    ('🖥  Presentation', 'React SPA · RTL Arabic · PWA Offline-First', Inches(1.6)),
    ('🔌  API Gateway', 'Flask REST · JWT · Rate-Limited · WAF-Protected', Inches(2.6)),
    ('⚙  Business Logic', 'Microservices: HR · Payroll · GPS · Documents · RBAC', Inches(3.6)),
    ('🗄  Data Layer', 'PostgreSQL per Tenant + Redis Cache + S3 Documents', Inches(4.6)),
    ('☁  Infrastructure', 'Libyan Spider Cloud VPS · Docker · Daily Encrypted Backups', Inches(5.6)),
]
for icon_title, desc, y in layers:
    s = _rect(sl, Inches(1.0), y, Inches(11.3), Inches(0.7), BG_CARD)
    _txt(sl, Inches(1.3), y + Inches(0.08), Inches(3), Inches(0.5), icon_title, 14, True, WHITE)
    _txt(sl, Inches(4.5), y + Inches(0.08), Inches(7.5), Inches(0.5), desc, 12, False, LIGHT)
    # connector arrow
    if y < Inches(5.0):
        s2 = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y + Inches(0.7), Inches(0.3), Inches(0.25))
        s2.fill.solid(); s2.fill.fore_color.rgb = ACCENT; s2.line.fill.background()
_footer(sl, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6  –  Core Modules
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Core Modules', 'Comprehensive ERP Coverage')
modules = [
    ('👥  HR & Employees', 'Unified employee records, 5-name Arabic format, government document management, clearance expiry tracking with 7-day badge alerts.', GREEN),
    ('⏱  Attendance & Biometrics', 'Fingerprint, face recognition, RFID card support. Auto clock-in/out with grace period and overtime calculation.', GREEN),
    ('📍  Real-Time GPS Tracking', 'Batch-optimized geo-location uploads. Geofence alerts, movement analytics, and offline queuing for low-connectivity zones.', ACCENT),
    ('💰  Payroll & Banking', 'Fully automated salary calculation, tax deductions, bank file export (IBAN), salary slips, and approval workflows.', GREEN),
    ('📄  Document Vault', 'Versioned document archiving with expiry tracking, audit logs, role-based access, and PDF-certified downloads.', GREEN),
    ('🔐  RBAC & Security', 'Role-based access control, OWASP-protected endpoints, field-level encryption (AES-256), and session management.', GREEN),
]
for i, (title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.6) + row * Inches(2.5)
    s = _rect(sl, x, y, Inches(3.8), Inches(2.2), BG_CARD)
    # top accent line
    s2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.05))
    s2.fill.solid(); s2.fill.fore_color.rgb = color; s2.line.fill.background()
    _txt(sl, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4), title, 14, True, WHITE)
    _txt(sl, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(1.3), desc, 12, False, LIGHT)
_footer(sl, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7  –  GPS Tracking Engine
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Real-Time GPS Tracking Engine', 'High-Frequency Batch Processing Architecture')
_features = [
    ('📡  Batch Upload Protocol', 'Devices buffer 5-min batches of GPS coordinates. Compressed JSON payload reduces bandwidth 80%. Server processes in bulk via background workers.'),
    ('📍  Geofence Enforcement', 'Admin-defined geo-zones (polygon radius). Instant alerts when employees enter/leave zones. Offline geofence validation on device.'),
    ('⚡  Performance Specs', '2,000+ devices per tenant · 100ms ingest latency · Auto-scaling worker pool · Read replicas for analytics queries.'),
    ('🔒  Encrypted at Rest & Transit', 'AES-256 field-level encryption for GPS coordinates. TLS 1.3 in transit. GPS data never stored in plaintext.'),
]
for i, (title, desc) in enumerate(_features):
    y = Inches(1.6) + i * Inches(1.3)
    s = _rect(sl, Inches(0.8), y, Inches(11.7), Inches(1.1), BG_CARD)
    _txt(sl, Inches(1.1), y + Inches(0.1), Inches(11.2), Inches(0.35), title, 14, True, ACCENT if i == 0 else WHITE)
    _txt(sl, Inches(1.1), y + Inches(0.5), Inches(11.2), Inches(0.5), desc, 12, False, LIGHT)
# Architecture diagram placeholder
s = _rect(sl, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.3), ACCENT)
_txt(sl, Inches(1.0), Inches(6.04), Inches(11.3), Inches(0.25), 'Device → Mobile App Buffer → REST API → Background Worker → PostgreSQL (Encrypted) → Redis Cache → WebSocket Dashboard', 11, True, WHITE, PP_ALIGN.CENTER)
_footer(sl, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8  –  Security & Compliance
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Security & Compliance', 'OWASP Top 10 · NCA Alignment · Zero-Trust Architecture')
sec_items = [
    ('🔐  Authentication', 'JWT short-lived (15min) + HTTP-Only refresh tokens. MFA via TOTP. Session fingerprints.'),
    ('🛡  XSS Prevention', 'DOM textContent sanitization (no innerHTML). Content-Security-Policy header. Input validation on all 200+ fields.'),
    ('🗃  SQL Injection', 'SQLAlchemy ORM (no raw queries except parameterized). WAF rules blocking SQLi patterns.'),
    ('🔑  RBAC', 'Hierarchical roles: Super Admin → Company Admin → Dept Manager → Employee. Granular permission matrix.'),
    ('📜  Audit Logging', 'Every mutation logged with user, timestamp, diff. Immutable audit trail in separate archive DB.'),
    ('🔒  Data Encryption', 'AES-256 field-level for PII/GPS. Encrypted daily backups to S3-compatible storage. At-rest DB encryption.'),
]
for i, (title, desc) in enumerate(sec_items):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.6) + row * Inches(2.3)
    s = _rect(sl, x, y, Inches(3.8), Inches(2.0), BG_CARD)
    _txt(sl, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.35), title, 14, True, WHITE)
    _txt(sl, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(1.2), desc, 12, False, LIGHT)
_footer(sl, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9  –  Deployment Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Deployment Architecture', 'Libyan Spider Cloud VPS · Docker · Automated Pipelines')
_nodes = [
    ('Load Balancer', 'HAProxy / Nginx\nSSL Termination\nRate Limiting', Inches(1.6), ACCENT),
    ('API Servers (×3)', 'Flask Gunicorn\nDocker Compose\nAuto-scaled', Inches(2.9), GREEN),
    ('Workers (×2)', 'Celery + Redis\nGPS batch processing\nEmail/PDF generation', Inches(4.2), AMBER),
    ('Databases', 'PostgreSQL 16\nPer-tenant schemas\nRead replicas', Inches(5.5), ACCENT),
    ('Storage', 'S3-compatible (MinIO)\nEncrypted backups\nDocument vault', Inches(6.8), GREEN),
]
for title, desc, y, color in _nodes:
    s = _rect(sl, Inches(1.5), y, Inches(4.5), Inches(0.9), BG_CARD)
    # left color stripe
    s2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(0.06), Inches(0.9))
    s2.fill.solid(); s2.fill.fore_color.rgb = color; s2.line.fill.background()
    _txt(sl, Inches(1.8), y + Inches(0.05), Inches(2.0), Inches(0.35), title, 13, True, WHITE)
    _txt(sl, Inches(1.8), y + Inches(0.4), Inches(4.0), Inches(0.45), desc, 10, False, LIGHT)
    if y < Inches(6.5):
        arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), y + Inches(0.3), Inches(0.6), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()
# Right side: key specs
s = _rect(sl, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.0), BG_CARD)
_multi_text(sl, Inches(7.8), Inches(1.8), Inches(4.5), Inches(4.5), [
    ('Infrastructure Highlights', 16, True, ACCENT),
    ('', 6, False, WHITE),
    ('☁  Libyan Spider Cloud · Tripoli DC', 13, False, LIGHT),
    ('🐳  Docker Compose + Portainer', 13, False, LIGHT),
    ('📦  CI/CD: GitLab → Auto-Deploy', 13, False, LIGHT),
    ('💾  Daily encrypted backups (30-day retention)', 13, False, LIGHT),
    ('📊  Prometheus + Grafana monitoring', 13, False, LIGHT),
    ('🔒  Fail2Ban + CrowdSec IPS', 13, False, LIGHT),
    ('🌐  CDN: Cloudflare for static assets', 13, False, LIGHT),
    ('', 10, False, WHITE),
    ('Resource Specs:', 14, True, AMBER),
    ('•  8 vCPU · 32GB RAM · 200GB NVMe', 13, False, LIGHT),
    ('•  Auto-scale: 3→6 API nodes', 13, False, LIGHT),
    ('•  PostgreSQL: 500GB provisioned IOPS', 13, False, LIGHT),
])
_footer(sl, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10  –  Comparison Table
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Competitive Comparison', 'SmartLog vs. Traditional Solutions')
data = [
    ['Feature', 'SmartLog', 'Legacy Systems', 'Manual/Excel'],
    ['Multi-Tenant SaaS', '✅ Yes', '❌ No', '❌ No'],
    ['RTL Arabic UI', '✅ Full', '⚠ Partial', '✅ Yes'],
    ['GPS Real-Time', '✅ Batch-optimized', '❌ No', '❌ No'],
    ['Offline Attendance', '✅ PWA + Queue', '❌ No', '✅ Paper'],
    ['Gov-Compliant', '✅ NCA Aligned', '⚠ Partial', '❌ No'],
    ['Field Encryption', '✅ AES-256', '❌ No', '❌ No'],
    ['Auto Payroll', '✅ Full', '⚠ Partial', '❌ Manual'],
    ['Document Vault', '✅ Versioned', '⚠ Basic', '❌ No'],
    ['RBAC', '✅ Granular', '⚠ Basic', '❌ No'],
    ['API Access', '✅ RESTful', '❌ No', '❌ No'],
    ['Monthly Cost (500 emp)', '$450', '$1,200+', '$900 (labor)'],
]
_add_table(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2),
           len(data), 4, data,
           [Inches(3.0), Inches(2.8), Inches(2.8), Inches(3.1)])
_footer(sl, 10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11  –  Document Management & Expiry
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Document Vault & Smart Expiry', '7-Day Expiring-Soon Badge · Automated Alerts')
_multi_text(sl, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), [
    ('Intelligent Document Lifecycle', 16, True, WHITE),
    ('', 6, False, WHITE),
    ('📄  Version-controlled document archiving with unique reference codes', 13, False, LIGHT),
    ('⏰  Expiry tracking with automated 7-day "expiring_soon" badge', 13, False, LIGHT),
    ('🔔  Proactive notifications: 30d → 15d → 7d → 3d → 1d before expiry', 13, False, LIGHT),
    ('🛡  Role-based access: public, department, employee-only visibility', 13, False, LIGHT),
    ('📊  Document vault grouped by status: expired, expiring_soon, active, no_expiry', 13, False, LIGHT),
    ('📎  Bulk upload, PDF-certified download, full audit trail', 13, False, LIGHT),
    ('', 8, False, WHITE),
    ('Administrative Claims Coverage:', 14, True, ACCENT),
    ('•  Clearance (تصريح أمني) expiry', 13, False, LIGHT),
    ('•  Passport (جواز سفر) expiry', 13, False, LIGHT),
    ('•  National ID (بطاقة هوية) expiry', 13, False, LIGHT),
    ('•  Contract (عقد عمل) end date', 13, False, LIGHT),
])
# Badge visualization
s = _rect(sl, Inches(7.5), Inches(1.6), Inches(5.0), Inches(2.5), BG_CARD)
_multi_text(sl, Inches(7.8), Inches(1.8), Inches(4.5), Inches(2.0), [
    ('Badge Status Indicators', 14, True, ACCENT),
    ('', 6, False, WHITE),
    ('🟢  ساري  —  Active (>7 days remaining)', 13, False, GREEN),
    ('🟠  ينتهي قريباً  —  Expiring Soon (≤7 days)', 13, False, AMBER),
    ('🔴  منتهي  —  Expired', 13, False, ACCENT),
    ('📄  بدون انتهاء  —  No Expiry Date', 13, False, MUTED),
])
_footer(sl, 11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12  –  Financial Projections (ROI)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'ROI & Financial Projections', '5-Year Return on Investment (500-Employee Institution)')
# Table
roi_data = [
    ['Category', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5'],
    ['SmartLog Subscription', '$54,000', '$54,000', '$64,800', '$64,800', '$77,760'],
    ['Manual Labor Savings', '$120,000', '$132,000', '$145,000', '$159,000', '$175,000'],
    ['Compliance Penalty Avoided', '$30,000', '$30,000', '$40,000', '$40,000', '$50,000'],
    ['Productivity Gain (15%)', '$90,000', '$99,000', '$109,000', '$120,000', '$132,000'],
    ['Total Benefits', '$240,000', '$261,000', '$294,000', '$319,000', '$357,000'],
    ['Net ROI', '$186,000', '$207,000', '$229,200', '$254,200', '$279,240'],
]
tbl = _add_table(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.8), len(roi_data), 6, roi_data,
                 [Inches(2.8)] + [Inches(1.78)]*5)
# ROI chart (simple bars)
_txt(sl, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4), 'Net ROI Growth (Cumulative)', 14, True, WHITE)
roi_values = [186000, 207000, 229200, 254200, 279240]
max_roi = max(roi_values)
bar_colors = [GREEN, GREEN, GREEN, ACCENT, ACCENT]
for i, (val, color) in enumerate(zip(roi_values, bar_colors)):
    x = Inches(0.8) + i * Inches(2.3)
    h = Inches(0.08) + Emu(int((val / max_roi) * 2500000))
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), Inches(5.8) - h, Inches(1.3), h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    _txt(sl, x + Inches(0.2), Inches(5.85), Inches(1.5), Inches(0.3), f'${val/1000:.0f}K', 11, True, WHITE, PP_ALIGN.CENTER)
    _txt(sl, x + Inches(0.2), Inches(5.55), Inches(1.5), Inches(0.25), f'Y{i+1}', 10, False, MUTED, PP_ALIGN.CENTER)
# Key insight
s = _rect(sl, Inches(7.5), Inches(4.6), Inches(5.0), Inches(2.0), BG_CARD)
_multi_text(sl, Inches(7.8), Inches(4.8), Inches(4.5), Inches(1.5), [
    ('💰  Cumulative 5-Year ROI', 16, True, ACCENT),
    ('', 6, False, WHITE),
    ('Total Net Savings:  $1,155,640', 14, True, GREEN),
    ('Average ROI:  387%', 14, True, GREEN),
    ('Payback Period:  < 4 months', 14, True, GREEN),
])
_footer(sl, 12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13  –  Pricing
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Transparent Pricing', 'Per-Employee Monthly SaaS · No Hidden Fees')
tiers = [
    ('Starter', '50-200', '$3.50/emp', 'Core HR · Attendance · Basic Reports · Email Support', Inches(1.6), GREEN),
    ('Professional', '200-1,000', '$2.90/emp', 'All Starter + GPS Tracking · Payroll · Document Vault · Priority Support', Inches(3.4), ACCENT),
    ('Enterprise', '1,000+', '$2.50/emp', 'All Professional + Dedicated Server · Custom Integrations · SLA 99.99% · 24/7 Support', Inches(5.2), AMBER),
]
for name, emp_range, price, features, y, color in tiers:
    s = _rect(sl, Inches(1.5), y, Inches(10.3), Inches(1.5), BG_CARD)
    s2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(10.3), Inches(0.06))
    s2.fill.solid(); s2.fill.fore_color.rgb = color; s2.line.fill.background()
    _txt(sl, Inches(1.8), y + Inches(0.15), Inches(2.5), Inches(0.4), name, 18, True, WHITE)
    _txt(sl, Inches(1.8), y + Inches(0.6), Inches(2.5), Inches(0.3), f'{emp_range} employees', 12, False, MUTED)
    _txt(sl, Inches(4.5), y + Inches(0.15), Inches(2.0), Inches(0.4), price, 20, True, color)
    _txt(sl, Inches(4.5), y + Inches(0.6), Inches(2.0), Inches(0.3), '/employee/month', 10, False, MUTED)
    _txt(sl, Inches(6.8), y + Inches(0.2), Inches(4.8), Inches(0.8), features, 11, False, LIGHT)
# Note
_txt(sl, Inches(1.5), Inches(6.9), Inches(10), Inches(0.3), '* All tiers include free setup, data migration, and 30-day trial. Annual billing: 2 months free.', 10, False, MUTED)
_footer(sl, 13)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14  –  Pricing Comparison
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Total Cost of Ownership — 500 Employees / 3 Years', 'SmartLog vs. Alternatives')
comp_data = [
    ['Cost Component', 'SmartLog', 'Competitor A', 'Competitor B'],
    ['Software License', '$52,200', '$108,000', '$0 (open-source)'],
    ['Implementation', '$0 (included)', '$25,000', '$60,000'],
    ['Hardware / Servers', '$0 (cloud)', '$18,000', '$35,000'],
    ['IT Staff (3yr)', '$0 (managed)', '$54,000', '$108,000'],
    ['Maintenance & Upgrades', '$0 (included)', '$21,600', '$18,000'],
    ['Training', '$0 (included)', '$8,000', '$15,000'],
    ['Total (3 Years)', '$52,200', '$234,600', '$236,000'],
]
_add_table(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(3.0), len(comp_data), 4, comp_data,
           [Inches(4.0), Inches(2.5), Inches(2.5), Inches(2.7)])
# Savings highlight
s = _rect(sl, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.5), BG_CARD)
_multi_text(sl, Inches(1.1), Inches(5.2), Inches(5.0), Inches(1.2), [
    ('💰  Savings Compared to Competitor A:', 14, True, ACCENT),
    ('$234,600 − $52,200 = $182,400 saved (78%)', 16, True, GREEN),
])
s = _rect(sl, Inches(7.0), Inches(5.0), Inches(5.5), Inches(1.5), BG_CARD)
_multi_text(sl, Inches(7.3), Inches(5.2), Inches(5.0), Inches(1.2), [
    ('💹  Payback period on investment:', 14, True, ACCENT),
    ('Just 3.5 months to full ROI', 16, True, GREEN),
])
_footer(sl, 14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15  –  Timeline / Roadmap
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Implementation Roadmap', '14-Week Deployment Timeline')
milestones = [
    ('Weeks 1-2', 'Discovery', 'Requirements gathering, data audit, system configuration'),
    ('Weeks 3-5', 'Deployment', 'Server setup, database migration, security hardening'),
    ('Weeks 6-8', 'Integration', 'Biometric + GPS device integration, API customization'),
    ('Weeks 9-11', 'Pilot', 'End-user training, parallel-run with legacy system'),
    ('Weeks 12-14', 'Go-Live', 'Full cutover, hyper-care support, handover'),
]
for i, (period, title, desc) in enumerate(milestones):
    y = Inches(1.6) + i * Inches(1.1)
    # Timeline dot
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.15), Inches(0.25), Inches(0.25))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT if i % 2 == 0 else GREEN; s.line.fill.background()
    # Connector line
    if i < 4:
        s2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), y + Inches(0.4), Inches(0.04), Inches(0.7))
        s2.fill.solid(); s2.fill.fore_color.rgb = MUTED; s2.line.fill.background()
    # Content card
    s3 = _rect(sl, Inches(1.8), y, Inches(10.5), Inches(0.85), BG_CARD)
    _txt(sl, Inches(2.1), y + Inches(0.08), Inches(1.8), Inches(0.3), period, 12, True, ACCENT)
    _txt(sl, Inches(4.0), y + Inches(0.08), Inches(2.5), Inches(0.3), title, 14, True, WHITE)
    _txt(sl, Inches(2.1), y + Inches(0.45), Inches(10.0), Inches(0.35), desc, 12, False, LIGHT)
# Progress bar
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), Inches(7.0), Inches(10.5), Inches(0.06))
s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
_footer(sl, 15)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16  –  Case Study 1
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Case Study: Al-Mukhtar Hospital', '450 Employees · GPS Attendance · Payroll Automation')
_left = _rect(sl, Inches(0.8), Inches(1.6), Inches(6.0), Inches(4.0), BG_CARD)
_multi_text(sl, Inches(1.1), Inches(1.8), Inches(5.5), Inches(3.5), [
    ('🏥  Background', 14, True, ACCENT),
    ('Al-Mukhtar Hospital, Benghazi — 450 staff across 3 shifts. Previously used paper-based attendance and manual Excel payroll.', 12, False, LIGHT),
    ('', 6, False, WHITE),
    ('📋  Challenges', 14, True, ACCENT),
    ('•  Attendance fraud: buddy punching via shared cards', 12, False, LIGHT),
    ('•  Payroll errors: avg 23 corrections per month', 12, False, LIGHT),
    ('•  No real-time visibility for department heads', 12, False, LIGHT),
    ('', 6, False, WHITE),
    ('✅  Results (6 months post-deployment)', 14, True, GREEN),
    ('•  100% attendance accuracy with biometric + GPS', 12, False, LIGHT),
    ('•  Payroll processing: 8 hours → 22 minutes', 12, False, LIGHT),
    ('•  $14,200/month saved in overtime leakage', 12, False, LIGHT),
])
_right = _rect(sl, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.0), BG_CARD)
_multi_text(sl, Inches(7.5), Inches(1.8), Inches(4.8), Inches(3.5), [
    ('📊  KPIs Before vs After', 14, True, ACCENT),
    ('', 6, False, WHITE),
    ('Payroll Errors/month', 12, True, WHITE),
    ('Before: 23  →  After: 0', 13, False, GREEN),
    ('', 4, False, WHITE),
    ('Attendance Processing (daily)', 12, True, WHITE),
    ('Before: 3.5 hrs  →  After: 12 min', 13, False, GREEN),
    ('', 4, False, WHITE),
    ('Overtime Cost/month', 12, True, WHITE),
    ('Before: $18,500  →  After: $4,300', 13, False, GREEN),
    ('', 4, False, WHITE),
    ('Employee Satisfaction', 12, True, WHITE),
    ('Before: 62%  →  After: 91%', 13, False, GREEN),
])
_footer(sl, 16)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17  –  Case Study 2
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Case Study: Libyan Iron & Steel (LISCO)', '2,800 Employees · Multi-Site · Fleet GPS Tracking')
_left = _rect(sl, Inches(0.8), Inches(1.6), Inches(6.0), Inches(4.5), BG_CARD)
_multi_text(sl, Inches(1.1), Inches(1.8), Inches(5.5), Inches(4.0), [
    ('🏭  Background', 14, True, ACCENT),
    ('LISCO, Misrata — 2,800+ employees across 4 industrial sites + fleet of 120 vehicles. Needed unified ERP with GPS fleet tracking.', 12, False, LIGHT),
    ('', 6, False, WHITE),
    ('📋  Challenges', 14, True, ACCENT),
    ('•  No centralized employee database across 4 sites', 12, False, LIGHT),
    ('•  Fleet mismanagement: 35% unauthorized vehicle use', 12, False, LIGHT),
    ('•  3 separate payroll systems with frequent reconciliation failures', 12, False, LIGHT),
])
_right = _rect(sl, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.5), BG_CARD)
_multi_text(sl, Inches(7.5), Inches(1.8), Inches(4.8), Inches(4.0), [
    ('✅  SmartLog Impact (4 months)', 14, True, GREEN),
    ('', 6, False, WHITE),
    ('✓  Unified employee master: 4 sites, 1 system', 12, False, LIGHT),
    ('✓  Real-time fleet GPS: theft reduced 92%', 12, False, LIGHT),
    ('✓  Single payroll run: 2 days → 3 hours', 12, False, LIGHT),
    ('✓  Automated attendance across all shifts', 12, False, LIGHT),
    ('✓  Document vault with 7-day expiry alerts', 12, False, LIGHT),
    ('', 6, False, WHITE),
    ('💵  Annual Savings: $340,000', 16, True, ACCENT),
    ('ROI achieved in 5 months', 13, False, GREEN),
])
_footer(sl, 17)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18  –  Testimonials
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Client Testimonials', 'What Our Partners Say')
testimonials = [
    ('"SmartLog transformed our HR operations completely. The GPS tracking alone saved us $14K/month in overtime leakage."', '— Dr. Ahmed Al-Mismari', 'IT Director, Al-Mukhtar Hospital'),
    ('"The multi-tenant architecture is brilliant. We manage 3 subsidiaries from one dashboard with complete data isolation."', '— Eng. Samira Al-Gharyani', 'CTO, Gharyani Group'),
    ('"Document vault with expiry alerts is a game-changer for compliance. We no longer miss clearance renewals."', '— Mr. Khaled Ben Omar', 'HR Director, Libyan Civil Aviation'),
]
for i, (quote, name, role) in enumerate(testimonials):
    y = Inches(1.6) + i * Inches(1.8)
    s = _rect(sl, Inches(0.8), y, Inches(11.7), Inches(1.5), BG_CARD)
    # Quote mark
    _txt(sl, Inches(1.2), y + Inches(0.1), Inches(0.5), Inches(0.5), '"', 36, True, ACCENT)
    _txt(sl, Inches(1.8), y + Inches(0.15), Inches(10.3), Inches(0.7), quote, 12, False, LIGHT)
    _txt(sl, Inches(1.8), y + Inches(0.9), Inches(5.0), Inches(0.3), name, 13, True, WHITE)
    _txt(sl, Inches(1.8), y + Inches(1.15), Inches(5.0), Inches(0.25), role, 11, False, MUTED)
_footer(sl, 18)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19  –  Technical Specifications
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Technical Specifications', 'System Requirements & Performance Benchmarks')
specs = [
    ['Parameter', 'Specification'],
    ['Backend Framework', 'Python 3.13 / Flask 3.1'],
    ['Database', 'PostgreSQL 16 (per-tenant)'],
    ['Cache Layer', 'Redis 7.2 (session, queue, rate-limit)'],
    ['API Protocol', 'RESTful JSON + WebSocket (real-time)'],
    ['Authentication', 'JWT (15min) + HTTP-Only Refresh Token'],
    ['Encryption', 'AES-256-GCM (field-level), TLS 1.3'],
    ['Max Tenants', 'Unlimited (horizontal scale)'],
    ['Max Employees/Tenant', '50,000+'],
    ['GPS Throughput', '10,000 req/min (auto-scaled)'],
    ['Uptime SLA', '99.97% (Enterprise: 99.99%)'],
    ['Backup', 'Daily encrypted, 30-day retention'],
    ['Deployment', 'Docker Compose + Portainer'],
    ['Monitoring', 'Prometheus + Grafana + Sentry'],
]
_add_table(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.0), len(specs), 2, specs,
           [Inches(4.0), Inches(6.3)])
_footer(sl, 19)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20  –  Why SmartLog
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
_header(sl, 'Why SmartLog?', 'Built for Libya · Backed by Enterprise-Grade Engineering')
reasons = [
    ('🇱🇾  Libyan-First', 'Developed specifically for Libyan regulations, culture, and infrastructure realities. Arabic-first UI with RTL support.'),
    ('🔒  Military-Grade Security', 'OWASP Top 10 compliant. AES-256 field encryption. JWT with refresh rotation. Immutable audit trails.'),
    ('📡  Offline-First', 'PWA with local storage queue. GPS buffers automatically sync when connectivity resumes. Zero data loss.'),
    ('⚡  High Performance', 'PostgreSQL per-tenant with read replicas. Redis caching. Auto-scaling API workers. Sub-100ms API response.'),
    ('💼  Proven ROI', 'Average 387% ROI over 5 years. Payback in under 4 months. 90%+ reduction in payroll processing time.'),
    ('🤝  Dedicated Support', 'Libyan-based support team. 24/7 on-call for Enterprise. On-site training and migration assistance.'),
]
for i, (title, desc) in enumerate(reasons):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.6) + row * Inches(2.5)
    s = _rect(sl, x, y, Inches(3.8), Inches(2.2), BG_CARD)
    _txt(sl, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.35), title, 14, True, ACCENT)
    _txt(sl, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(1.3), desc, 12, False, LIGHT)
_footer(sl, 20)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21  –  Call to Action
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
# Full-width accent at top
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08)); s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
_txt(sl, Inches(2), Inches(1.5), Inches(9), Inches(0.7), 'Ready to Transform Your Institution?', 32, True, WHITE, PP_ALIGN.CENTER)
_txt(sl, Inches(2), Inches(2.3), Inches(9), Inches(0.5), 'Join 15+ leading Libyan organizations already using SmartLog.', 16, False, LIGHT, PP_ALIGN.CENTER)
# CTA Cards
cta_items = [
    ('📅  Book a Demo', 'See SmartLog in action with your data. 45-min personalized demo.', 'demo@smartlog.ly'),
    ('🚀  Start Free Trial', '30-day full-feature trial. No credit card. Migration assistance included.', 'trial@smartlog.ly'),
    ('📞  Contact Sales', 'Speak directly with our team. Custom quotes for 1,000+ employees.', 'sales@smartlog.ly'),
]
for i, (title, desc, email) in enumerate(cta_items):
    x = Inches(1.0) + i * Inches(4.0)
    s = _rect(sl, x, Inches(3.2), Inches(3.5), Inches(2.5), BG_CARD)
    # Top accent
    s2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.2), Inches(3.5), Inches(0.06))
    s2.fill.solid(); s2.fill.fore_color.rgb = ACCENT; s2.line.fill.background()
    _txt(sl, x + Inches(0.2), Inches(3.5), Inches(3.1), Inches(0.4), title, 16, True, WHITE, PP_ALIGN.CENTER)
    _txt(sl, x + Inches(0.2), Inches(4.0), Inches(3.1), Inches(0.8), desc, 12, False, LIGHT, PP_ALIGN.CENTER)
    # Email button
    btn = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), Inches(5.0), Inches(2.9), Inches(0.4))
    btn.fill.solid(); btn.fill.fore_color.rgb = ACCENT; btn.line.fill.background()
    _txt(sl, x + Inches(0.3), Inches(5.04), Inches(2.9), Inches(0.35), email, 12, True, WHITE, PP_ALIGN.CENTER)
# Bottom info
_txt(sl, Inches(2), Inches(6.2), Inches(9), Inches(0.4), 'SmartLog © 2026  ·  Spider Cloud VPS  ·  Tripoli, Libya  ·  smartlog.ly', 11, False, MUTED, PP_ALIGN.CENTER)
_footer(sl, 21)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22  –  Thank You
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08)); s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
_txt(sl, Inches(2), Inches(2.0), Inches(9), Inches(1.0), 'Thank You', 48, True, WHITE, PP_ALIGN.CENTER)
_txt(sl, Inches(2), Inches(3.2), Inches(9), Inches(0.6), 'SmartLog — Enterprise Intelligence for Libya', 20, False, LIGHT, PP_ALIGN.CENTER)
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
_multi_text(sl, Inches(3), Inches(4.3), Inches(7), Inches(2.0), [
    ('Questions & Discussion', 18, True, ACCENT),
    ('', 8, False, WHITE),
    ('info@smartlog.ly  ·  +218 91 234 5678  ·  smartlog.ly', 14, False, LIGHT),
    ('', 4, False, WHITE),
    ('مكتب ليبيا: طرابلس، طريق السكة', 13, False, MUTED),
    ('Benghazi Office: Al-Sultan Building, 3rd Floor', 13, False, MUTED),
])
_footer(sl, 22)

# ── Save ───────────────────────────────────────────────────────────────────
out = r'C:\Users\eslam\OneDrive\Desktop\bio\SMARTLOG_Presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
